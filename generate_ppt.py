from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    def add_slide(title_text, points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.word_wrap = True
        
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.space_after = Pt(10)

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "HireByte AI"
    subtitle.text = "The Intelligent Interview Twin\nPersonalized Mock Interviews with Real-Time Vision Analytics"

    # Slide 2: Proposed Solution (Problem & Fix)
    add_slide("Proposed Solution", [
        "Problem: Interview preparation is often generic, expensive, or lacks real-time feedback on body language.",
        "Solution: HireByte AI provides an end-to-end simulated environment that mimics real professional interviews.",
        "How it works: It reads YOUR resume and the specific Job Description to generate relevant, high-quality questions.",
        "The Edge: Provides live 'Vision Analytics' to track your focus, confidence, and stress levels while you speak."
    ])

    # Slide 3: Technical Approach - Tech Stack
    add_slide("Technical Approach: Tech Stack", [
        "Frontend: React with TypeScript for a smooth, interactive dashboard.",
        "Backend: Python (FastAPI) for high-speed data processing.",
        "AI Brain: GPT-4o for intelligent conversation and personalized feedback.",
        "Voice: Whisper (Speech-to-Text) and OpenAI TTS (AI Speaking).",
        "Vision: OpenCV for facial tracking and focus measurement."
    ])

    # Slide 4: Technical Approach - Implementation
    add_slide("Technical Approach: Implementation", [
        "Context Mapping: Seamlessly combines Resume details with Job requirements.",
        "Dual-Stream System: Processes Voice and Video simultaneously via WebSockets.",
        "Real-Time Stream: Uses 10 FPS capture for lag-free posture and eye-contact analysis.",
        "Data Dashboard: Displays time-series charts of your performance immediately after the interview."
    ])

    # Slide 5: Feasibility & Viability
    add_slide("Feasibility and Viability", [
        "Feasibility: Built using browser-native APIs; no expensive hardware required.",
        "Potential Risks: Background lighting or network latency can affect real-time metrics.",
        "Strategies: Implemented 'Adaptive Smoothing' algorithms to filter noise and ensure stable, accurate data even in normal room lighting."
    ])

    # Slide 6: Impact and Benefits
    add_slide("Impact and Benefits", [
        "Impact: Helps students and job seekers overcome interview anxiety in a low-risk environment.",
        "Economic: Reduces the time spent in the job search by improving first-round performance.",
        "Social: Levels the playing field for introverts or non-native speakers by providing honest, AI-driven coaching.",
        "Future-Ready: High potential for use in college placement cells and corporate HR screening."
    ])

    # Save the presentation
    output_path = "HireByte_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
